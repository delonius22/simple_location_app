"""
Microsoft File Converter for RSA-TOTP Encryption System.

This module handles conversion of Microsoft Office files (Word, Excel, PowerPoint)
to JSON format for encryption, and back to their original format for decryption.
"""

import os
import json
import base64
from pathlib import Path
from typing import Dict, Any, Tuple, Optional, Union

try:
    # For Excel files
    import pandas as pd
    import openpyxl
    
    # For Word documents
    import docx
    from docx.document import Document as DocxDocument
    
    # For PowerPoint
    import pptx
    from pptx import Presentation
    
    # Optional - for more advanced Word document processing
    import mammoth
    
except ImportError:
    print("Warning: Some Microsoft Office file conversion libraries are missing.")
    print("Please install them using: pip install pandas openpyxl python-docx python-pptx mammoth")


class MSFileConverter:
    """Converts Microsoft Office files to and from JSON format for encryption."""
    
    def __init__(self):
        """Initialize the file converter."""
        self.supported_extensions = {
            # Microsoft Word
            '.docx': self._convert_word_to_json,
            '.doc': self._convert_word_to_json,
            
            # Microsoft Excel
            '.xlsx': self._convert_excel_to_json,
            '.xls': self._convert_excel_to_json,
            '.xlsm': self._convert_excel_to_json,
            
            # Microsoft PowerPoint
            '.pptx': self._convert_powerpoint_to_json,
            '.ppt': self._convert_powerpoint_to_json,
            
            # Already JSON
            '.json': self._load_json
        }
        
        self.reverse_converters = {
            'word': self._convert_json_to_word,
            'excel': self._convert_json_to_excel,
            'powerpoint': self._convert_json_to_powerpoint,
            'json': self._save_json
        }
    
    def convert_to_json(self, file_path: str) -> Tuple[Dict[str, Any], str]:
        """
        Convert a Microsoft Office file to JSON.
        
        Args:
            file_path: Path to the Microsoft Office file
            
        Returns:
            Tuple of (json_data, file_type)
            
        Raises:
            ValueError: If the file format is not supported
        """
        file_path = Path(file_path)
        extension = file_path.suffix.lower()
        
        if extension not in self.supported_extensions:
            raise ValueError(f"Unsupported file format: {extension}")
        
        # Call the appropriate converter based on file extension
        converter = self.supported_extensions[extension]
        return converter(file_path)
    
    def convert_from_json(self, json_data: Dict[str, Any], output_path: str) -> str:
        """
        Convert JSON back to the original Microsoft Office format.
        
        Args:
            json_data: The JSON data to convert
            output_path: Path where to save the converted file
            
        Returns:
            Path to the converted file
            
        Raises:
            ValueError: If the file format is not supported or if the JSON lacks required metadata
        """
        if 'file_type' not in json_data.get('metadata', {}):
            raise ValueError("JSON data does not contain file type information")
        
        file_type = json_data['metadata']['file_type']
        
        if file_type not in self.reverse_converters:
            raise ValueError(f"Unsupported conversion from JSON to {file_type}")
        
        # Call the appropriate converter
        converter = self.reverse_converters[file_type]
        return converter(json_data, output_path)
    
    def is_supported_format(self, file_path: str) -> bool:
        """
        Check if the file format is supported.
        
        Args:
            file_path: Path to the file
            
        Returns:
            True if the format is supported, False otherwise
        """
        extension = Path(file_path).suffix.lower()
        return extension in self.supported_extensions
    
    def _load_json(self, file_path: Path) -> Tuple[Dict[str, Any], str]:
        """Load a JSON file directly."""
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        return data, 'json'
    
    def _save_json(self, json_data: Dict[str, Any], output_path: str) -> str:
        """Save JSON data to a file."""
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(json_data.get('content', {}), f, indent=2)
        
        return output_path
    
    def _convert_word_to_json(self, file_path: Path) -> Tuple[Dict[str, Any], str]:
        """Convert a Word document to JSON."""
        try:
            # Create a structure to hold the document content
            doc_json = {
                'metadata': {
                    'file_type': 'word',
                    'original_filename': file_path.name,
                    'extension': file_path.suffix.lower()
                },
                'content': {
                    'paragraphs': [],
                    'tables': []
                },
                'binary_content': ''
            }
            
            # Save binary content for exact reconstruction
            with open(file_path, 'rb') as f:
                binary_content = f.read()
                doc_json['binary_content'] = base64.b64encode(binary_content).decode('utf-8')
            
            # Extract text content
            try:
                doc = docx.Document(file_path)
                
                # Extract paragraphs
                for para in doc.paragraphs:
                    if para.text.strip():  # Skip empty paragraphs
                        doc_json['content']['paragraphs'].append({
                            'text': para.text,
                            'style': para.style.name if para.style else 'Normal'
                        })
                
                # Extract tables
                for i, table in enumerate(doc.tables):
                    table_data = []
                    for row in table.rows:
                        row_data = []
                        for cell in row.cells:
                            row_data.append(cell.text)
                        table_data.append(row_data)
                    
                    doc_json['content']['tables'].append(table_data)
                
            except Exception as e:
                # If python-docx fails, try using mammoth for text extraction
                try:
                    with open(file_path, 'rb') as docx_file:
                        result = mammoth.extract_raw_text(docx_file)
                        doc_json['content']['paragraphs'].append({
                            'text': result.value,
                            'style': 'Normal'
                        })
                except:
                    # If all else fails, at least we have the binary content
                    doc_json['content']['paragraphs'].append({
                        'text': f"Error extracting text: {str(e)}",
                        'style': 'Error'
                    })
            
            return doc_json, 'word'
            
        except Exception as e:
            raise ValueError(f"Error converting Word document: {str(e)}")
    
    def _convert_json_to_word(self, json_data: Dict[str, Any], output_path: str) -> str:
        """Convert JSON back to a Word document."""
        # The most reliable way to reconstruct the Word document is to use the binary content
        if 'binary_content' in json_data and json_data['binary_content']:
            try:
                binary_data = base64.b64decode(json_data['binary_content'])
                with open(output_path, 'wb') as f:
                    f.write(binary_data)
                return output_path
            except Exception as e:
                raise ValueError(f"Error reconstructing Word document from binary: {str(e)}")
        
        # Fallback: Create a new document from the text content
        try:
            doc = docx.Document()
            
            # Add paragraphs
            for para in json_data.get('content', {}).get('paragraphs', []):
                p = doc.add_paragraph(para.get('text', ''))
                if 'style' in para and para['style'] != 'Error':
                    try:
                        p.style = para['style']
                    except:
                        pass  # Style might not exist in the template
            
            # Add tables
            for table_data in json_data.get('content', {}).get('tables', []):
                if not table_data:
                    continue
                    
                table = doc.add_table(rows=len(table_data), cols=len(table_data[0]))
                
                for i, row_data in enumerate(table_data):
                    for j, cell_text in enumerate(row_data):
                        if i < len(table.rows) and j < len(table.rows[i].cells):
                            table.rows[i].cells[j].text = cell_text
            
            # Save the document
            doc.save(output_path)
            return output_path
            
        except Exception as e:
            raise ValueError(f"Error creating Word document: {str(e)}")
    
    def _convert_excel_to_json(self, file_path: Path) -> Tuple[Dict[str, Any], str]:
        """Convert an Excel file to JSON."""
        try:
            # Create a structure to hold the Excel content
            excel_json = {
                'metadata': {
                    'file_type': 'excel',
                    'original_filename': file_path.name,
                    'extension': file_path.suffix.lower()
                },
                'content': {
                    'sheets': {}
                },
                'binary_content': ''
            }
            
            # Save binary content for exact reconstruction
            with open(file_path, 'rb') as f:
                binary_content = f.read()
                excel_json['binary_content'] = base64.b64encode(binary_content).decode('utf-8')
            
            # Use pandas to read Excel sheets
            try:
                # Get sheet names
                xlsx = pd.ExcelFile(file_path)
                sheet_names = xlsx.sheet_names
                
                # Convert each sheet to JSON
                for sheet_name in sheet_names:
                    df = pd.read_excel(file_path, sheet_name=sheet_name)
                    excel_json['content']['sheets'][sheet_name] = json.loads(df.to_json(orient='records'))
                    
                # Add sheet structure information
                workbook = openpyxl.load_workbook(file_path, data_only=True)
                excel_json['metadata']['sheet_names'] = sheet_names
                
            except Exception as e:
                # If pandas fails, at least we have the binary content
                excel_json['content']['error'] = f"Error extracting Excel data: {str(e)}"
            
            return excel_json, 'excel'
            
        except Exception as e:
            raise ValueError(f"Error converting Excel file: {str(e)}")
    
    def _convert_json_to_excel(self, json_data: Dict[str, Any], output_path: str) -> str:
        """Convert JSON back to an Excel file."""
        # The most reliable way to reconstruct the Excel file is to use the binary content
        if 'binary_content' in json_data and json_data['binary_content']:
            try:
                binary_data = base64.b64decode(json_data['binary_content'])
                with open(output_path, 'wb') as f:
                    f.write(binary_data)
                return output_path
            except Exception as e:
                print(f"Warning: Failed to reconstruct Excel from binary: {str(e)}")
                print("Falling back to data-only reconstruction...")
        
        # Fallback: Create a new Excel file from the data
        try:
            with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
                for sheet_name, sheet_data in json_data.get('content', {}).get('sheets', {}).items():
                    # Convert sheet data to DataFrame
                    df = pd.DataFrame(sheet_data)
                    # Write to Excel
                    df.to_excel(writer, sheet_name=sheet_name, index=False)
            
            return output_path
            
        except Exception as e:
            raise ValueError(f"Error creating Excel file: {str(e)}")
    
    def _convert_powerpoint_to_json(self, file_path: Path) -> Tuple[Dict[str, Any], str]:
        """Convert a PowerPoint file to JSON."""
        try:
            # Create a structure to hold the PowerPoint content
            ppt_json = {
                'metadata': {
                    'file_type': 'powerpoint',
                    'original_filename': file_path.name,
                    'extension': file_path.suffix.lower()
                },
                'content': {
                    'slides': []
                },
                'binary_content': ''
            }
            
            # Save binary content for exact reconstruction
            with open(file_path, 'rb') as f:
                binary_content = f.read()
                ppt_json['binary_content'] = base64.b64encode(binary_content).decode('utf-8')
            
            # Extract text content
            try:
                pres = Presentation(file_path)
                
                # Get slide count
                ppt_json['metadata']['slide_count'] = len(pres.slides)
                
                # Extract text from each slide
                for i, slide in enumerate(pres.slides):
                    slide_data = {
                        'slide_number': i + 1,
                        'slide_id': str(i),  # Placeholder
                        'shapes': []
                    }
                    
                    # Extract text from shapes
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_data['shapes'].append({
                                'text': shape.text,
                                'type': 'text'
                            })
                    
                    ppt_json['content']['slides'].append(slide_data)
                
            except Exception as e:
                # If python-pptx fails, at least we have the binary content
                ppt_json['content']['error'] = f"Error extracting PowerPoint data: {str(e)}"
            
            return ppt_json, 'powerpoint'
            
        except Exception as e:
            raise ValueError(f"Error converting PowerPoint file: {str(e)}")
    
    def _convert_json_to_powerpoint(self, json_data: Dict[str, Any], output_path: str) -> str:
        """Convert JSON back to a PowerPoint file."""
        # The most reliable way to reconstruct the PowerPoint is to use the binary content
        if 'binary_content' in json_data and json_data['binary_content']:
            try:
                binary_data = base64.b64decode(json_data['binary_content'])
                with open(output_path, 'wb') as f:
                    f.write(binary_data)
                return output_path
            except Exception as e:
                print(f"Warning: Failed to reconstruct PowerPoint from binary: {str(e)}")
                print("Falling back to content-only reconstruction...")
        
        # Fallback: Create a simple PowerPoint with text content
        try:
            pres = Presentation()
            
            # Add slides with text
            for slide_data in json_data.get('content', {}).get('slides', []):
                slide = pres.slides.add_slide(pres.slide_layouts[5])  # Title and Content layout
                
                # Add text from shapes
                for i, shape_data in enumerate(slide_data.get('shapes', [])):
                    if shape_data.get('type') == 'text' and shape_data.get('text'):
                        textbox = slide.shapes.add_textbox(
                            left=100, 
                            top=100 + i * 50, 
                            width=400, 
                            height=50
                        )
                        textbox.text_frame.text = shape_data['text']
            
            # Save the presentation
            pres.save(output_path)
            return output_path
            
        except Exception as e:
            raise ValueError(f"Error creating PowerPoint file: {str(e)}")


# Self-test code
if __name__ == "__main__":
    import sys
    
    if len(sys.argv) < 2:
        print("Usage: python file_converter.py <file_path>")
        sys.exit(1)
    
    file_path = sys.argv[1]
    
    try:
        converter = MSFileConverter()
        
        if not converter.is_supported_format(file_path):
            print(f"Unsupported file format: {Path(file_path).suffix}")
            sys.exit(1)
        
        print(f"Converting {file_path} to JSON...")
        json_data, file_type = converter.convert_to_json(file_path)
        
        # Save the JSON for inspection
        json_path = f"{Path(file_path).stem}_converted.json"
        with open(json_path, 'w', encoding='utf-8') as f:
            json.dump(json_data, f, indent=2)
        
        print(f"Saved JSON to {json_path}")
        
        # Test conversion back to original format
        output_path = f"{Path(file_path).stem}_reconstructed{Path(file_path).suffix}"
        converter.convert_from_json(json_data, output_path)
        
        print(f"Converted back to {output_path}")
        
    except Exception as e:
        print(f"Error: {str(e)}")
        sys.exit(1)